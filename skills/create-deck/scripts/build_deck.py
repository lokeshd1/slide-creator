#!/usr/bin/env python3
from __future__ import annotations

import argparse
from copy import deepcopy
from dataclasses import dataclass, field
import json
from pathlib import Path
import re
import sys

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None


@dataclass
class Block:
    type: str
    text: str = ""
    items: list[str] = field(default_factory=list)
    level: int = 1
    src: str | None = None
    alt: str = ""
    data: dict[str, object] = field(default_factory=dict)


@dataclass
class Slide:
    title: str = ""
    subtitle: str = ""
    blocks: list[Block] = field(default_factory=list)
    notes: str = ""
    layout: str = "content"


@dataclass
class Deck:
    title: str = "Untitled Deck"
    slides: list[Slide] = field(default_factory=list)
    source_path: Path | None = None
    metadata: dict[str, str] = field(default_factory=dict)


@dataclass(frozen=True)
class Theme:
    heading_font: str
    body_font: str
    background: tuple[int, int, int]
    title_background: tuple[int, int, int]
    surface: tuple[int, int, int]
    title: tuple[int, int, int]
    body: tuple[int, int, int]
    muted: tuple[int, int, int]
    accent: tuple[int, int, int]
    accent_2: tuple[int, int, int]
    accent_3: tuple[int, int, int]
    accent_4: tuple[int, int, int]
    on_accent: tuple[int, int, int]


THEMES = {
    "modern": Theme("Aptos Display", "Aptos", (247, 249, 252), (20, 28, 42), (255, 255, 255), (28, 38, 54), (54, 65, 83), (105, 116, 132), (34, 112, 177), (22, 163, 141), (235, 142, 54), (126, 87, 194), (255, 255, 255)),
    "executive": Theme("Georgia", "Aptos", (250, 249, 246), (36, 36, 33), (255, 255, 255), (35, 35, 32), (64, 62, 58), (126, 119, 108), (158, 111, 45), (72, 104, 119), (121, 87, 66), (88, 96, 76), (255, 255, 255)),
    "studio": Theme("Aptos Display", "Aptos", (248, 247, 252), (49, 32, 76), (255, 255, 255), (43, 36, 56), (65, 58, 76), (116, 105, 131), (123, 82, 193), (216, 85, 124), (41, 166, 185), (245, 168, 66), (255, 255, 255)),
    "mono": Theme("Helvetica", "Helvetica", (250, 250, 249), (24, 24, 27), (255, 255, 255), (24, 24, 27), (63, 63, 70), (113, 113, 122), (24, 24, 27), (82, 82, 91), (161, 161, 170), (212, 212, 216), (255, 255, 255)),
}

SLIDE_BREAK = re.compile(r"^\s*---+\s*$")
HEADING = re.compile(r"^(#{1,6})\s+(.+)$")
BULLET = re.compile(r"^\s*[-*+]\s+(.+)$")
NUMBERED = re.compile(r"^\s*\d+[.)]\s+(.+)$")
IMAGE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
METRIC = re.compile(r"(?<!\w)([$]?\d+(?:[.,]\d+)?\s?(?:%|x|k|m|b|K|M|B)?)(?!\w)")


def main() -> int:
    parser = argparse.ArgumentParser(description="Build or inspect a themed PPTX deck from Markdown.")
    parser.add_argument("source", type=Path, help="Markdown or text source file.")
    parser.add_argument("--output", "-o", type=Path, help="Output PPTX path.")
    parser.add_argument("--theme", default="modern", choices=sorted(THEMES), help="Visual theme.")
    parser.add_argument("--no-infographics", action="store_true", help="Disable automatic infographics.")
    parser.add_argument("--inspect", action="store_true", help="Print parsed deck JSON instead of building.")
    parser.add_argument("--infographics", action="store_true", help="Show generated infographics in inspect output.")
    args = parser.parse_args()

    deck = load_deck(args.source)
    if args.infographics or (not args.inspect and not args.no_infographics):
        deck = add_auto_infographics(deck)

    if args.inspect:
        print(json.dumps(deck_to_dict(deck), indent=2))
        return 0

    if not args.output:
        parser.error("--output is required unless --inspect is used")

    if Presentation is None:
        raise RuntimeError("python-pptx is required. Install it with: pip install python-pptx")

    output = PptxRenderer(THEMES[args.theme]).render(deck, args.output.expanduser().resolve())
    print(output)
    return 0


def load_deck(source_path: Path) -> Deck:
    source_path = source_path.expanduser().resolve()
    if source_path.suffix.lower() == ".docx":
        markdown = read_docx(source_path)
    else:
        markdown = source_path.read_text(encoding="utf-8")
    return parse_markdown(markdown, source_path)


def read_docx(source_path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("Reading .docx files requires python-docx: pip install python-docx") from exc

    lines = []
    for paragraph in Document(source_path).paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name.lower() if paragraph.style else ""
        if "title" in style:
            lines.append(f"# {text}")
        elif "heading 1" in style:
            lines.append(f"---\n# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "list" in style:
            lines.append(f"- {text}")
        else:
            lines.append(text)
    return "\n\n".join(lines)


def parse_markdown(markdown: str, source_path: Path | None = None) -> Deck:
    metadata, body = split_front_matter(markdown)
    slides = [
        parse_slide(chunk, is_title_slide=index == 0)
        for index, chunk in enumerate(split_slides(body))
        if chunk.strip()
    ]
    return Deck(metadata.get("title") or deck_title(slides), slides, source_path, metadata)


def split_front_matter(markdown: str) -> tuple[dict[str, str], str]:
    lines = markdown.splitlines()
    if not lines or lines[0].strip() != "---":
        return {}, markdown

    metadata = {}
    for index, line in enumerate(lines[1:], start=1):
        if line.strip() == "---":
            return metadata, "\n".join(lines[index + 1 :])
        if ":" in line:
            key, value = line.split(":", 1)
            metadata[key.strip().lower()] = value.strip()
    return {}, markdown


def split_slides(markdown: str) -> list[str]:
    slides = [[]]
    in_code = False
    for line in markdown.splitlines():
        if line.strip().startswith("```"):
            in_code = not in_code
        if not in_code and SLIDE_BREAK.match(line):
            slides.append([])
        else:
            slides[-1].append(line)
    return ["\n".join(slide).strip() for slide in slides]


def parse_slide(markdown: str, is_title_slide: bool = False) -> Slide:
    content, notes = split_notes(markdown)
    slide = Slide(notes=notes)
    paragraph_lines, bullets, numbered, code_lines = [], [], [], []
    in_code = False

    def flush_paragraph() -> None:
        if paragraph_lines:
            text = " ".join(paragraph_lines).strip()
            if is_title_slide and not slide.subtitle and slide.title and not slide.blocks:
                slide.subtitle = text
            else:
                slide.blocks.append(Block("paragraph", text=text))
            paragraph_lines.clear()

    def flush_bullets() -> None:
        if bullets:
            slide.blocks.append(Block("bullets", items=list(bullets)))
            bullets.clear()

    def flush_numbered() -> None:
        if numbered:
            slide.blocks.append(Block("numbered", items=list(numbered)))
            numbered.clear()

    for raw_line in content.splitlines():
        line = raw_line.rstrip()
        if line.strip().startswith("```"):
            if in_code:
                slide.blocks.append(Block("code", text="\n".join(code_lines)))
                code_lines.clear()
            else:
                flush_paragraph()
                flush_bullets()
                flush_numbered()
            in_code = not in_code
            continue
        if in_code:
            code_lines.append(raw_line)
            continue
        if not line.strip():
            flush_paragraph()
            flush_bullets()
            flush_numbered()
            continue

        image = IMAGE.search(line)
        if image:
            flush_paragraph()
            flush_bullets()
            flush_numbered()
            slide.blocks.append(Block("image", alt=image.group(1), src=image.group(2)))
            continue

        heading = HEADING.match(line)
        if heading:
            flush_paragraph()
            flush_bullets()
            flush_numbered()
            text = heading.group(2).strip()
            if not slide.title:
                slide.title = text
            else:
                slide.blocks.append(Block("heading", text=text, level=len(heading.group(1))))
            continue

        bullet = BULLET.match(line)
        if bullet:
            flush_paragraph()
            flush_numbered()
            bullets.append(bullet.group(1).strip())
            continue

        ordered = NUMBERED.match(line)
        if ordered:
            flush_paragraph()
            flush_bullets()
            numbered.append(ordered.group(1).strip())
            continue

        paragraph_lines.append(line.strip())

    if in_code and code_lines:
        slide.blocks.append(Block("code", text="\n".join(code_lines)))
    flush_paragraph()
    flush_bullets()
    flush_numbered()
    return slide


def split_notes(markdown: str) -> tuple[str, str]:
    lines = markdown.splitlines()
    for index, line in enumerate(lines):
        if line.strip().lower() in {"notes:", "speaker notes:"}:
            return "\n".join(lines[:index]).strip(), "\n".join(lines[index + 1 :]).strip()
    return markdown, ""


def deck_title(slides: list[Slide]) -> str:
    return next((slide.title for slide in slides if slide.title), "Untitled Deck")


def add_auto_infographics(deck: Deck) -> Deck:
    enriched = deepcopy(deck)
    for index, slide in enumerate(enriched.slides):
        if index == 0 or any(block.type == "infographic" for block in slide.blocks):
            continue
        infographic = infer_infographic(slide)
        if infographic:
            slide.blocks.append(infographic)
    return enriched


def infer_infographic(slide: Slide) -> Block | None:
    metrics = extract_metrics(slide)
    if metrics:
        return Block("infographic", data={"variant": "metrics", "title": "Key Metrics", "metrics": metrics[:4]})
    numbered = first_list(slide, "numbered")
    if 2 <= len(numbered) <= 5:
        return Block("infographic", data={"variant": "process", "title": "Flow", "steps": numbered})
    bullets = first_list(slide, "bullets")
    if 2 <= len(bullets) <= 4:
        return Block("infographic", data={"variant": "highlights", "title": "Highlights", "items": bullets})
    return None


def extract_metrics(slide: Slide) -> list[dict[str, str]]:
    metrics = []
    for text in slide_text_units(slide):
        for match in METRIC.finditer(text):
            value = match.group(1).strip()
            label = metric_label(text, value)
            if not any(metric["value"] == value and metric["label"] == label for metric in metrics):
                metrics.append({"value": value, "label": label})
    return metrics


def slide_text_units(slide: Slide) -> list[str]:
    units = [slide.subtitle] if slide.subtitle else []
    for block in slide.blocks:
        if block.type in {"paragraph", "heading", "code"} and block.text:
            units.append(block.text)
        elif block.type in {"bullets", "numbered"}:
            units.extend(block.items)
    return units


def first_list(slide: Slide, block_type: str) -> list[str]:
    return next((block.items for block in slide.blocks if block.type == block_type), [])


def metric_label(text: str, value: str) -> str:
    for_pattern = re.search(rf"(.+?)\s+for\s+{re.escape(value)}\s+of\s+(.+)", text, flags=re.IGNORECASE)
    if for_pattern:
        label = f"{for_pattern.group(1).strip(' .:-')} ({for_pattern.group(2).strip(' .:-')})"
    else:
        label = text.replace(value, "").strip(" .:-")
    return (label[:55].rstrip() + "...") if len(label) > 58 else label or "Metric"


class PptxRenderer:
    def __init__(self, theme: Theme) -> None:
        self.theme = theme

    def render(self, deck: Deck, output_path: Path) -> Path:
        presentation = Presentation()
        presentation.core_properties.title = deck.title
        image_base_path = deck.source_path.parent if deck.source_path else Path.cwd()
        for index, slide_model in enumerate(deck.slides):
            if index == 0 and slide_model.subtitle:
                self.add_title_slide(presentation, slide_model)
            else:
                self.add_content_slide(presentation, slide_model, image_base_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        return output_path

    def add_title_slide(self, presentation, slide_model: Slide) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.set_background(slide, self.theme.title_background)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = rgb(self.theme.accent_2)
        accent.line.fill.background()
        self.add_text(slide, slide_model.title, 0.82, 2.1, 8.7, 1.2, 44, self.theme.on_accent, bold=True, font=self.theme.heading_font)
        self.add_text(slide, slide_model.subtitle, 0.86, 3.35, 7.7, 0.7, 20, tint(self.theme.on_accent, self.theme.title_background, 0.2), font=self.theme.body_font)
        self.add_notes(slide, slide_model.notes)

    def add_content_slide(self, presentation, slide_model: Slide, image_base_path: Path) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.set_background(slide, self.theme.background)
        self.add_header(slide, slide_model.title or "Untitled")
        infographic = next((block for block in slide_model.blocks if block.type == "infographic"), None)
        width = 4.65 if infographic else 8.7
        self.add_panel(slide, 0.76, 1.45, width, 5.3)
        box = slide.shapes.add_textbox(Inches(0.76), Inches(1.45), Inches(width), Inches(5.3))
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.28)
        frame.margin_right = Inches(0.25)
        frame.margin_top = Inches(0.22)
        frame.clear()
        first = True
        for block in slide_model.blocks:
            if block.type == "infographic":
                continue
            if block.type == "image":
                self.add_image(slide, block, image_base_path)
                continue
            paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            self.write_block(paragraph, block)
        if infographic:
            self.add_infographic(slide, infographic)
        self.add_notes(slide, slide_model.notes)

    def write_block(self, paragraph, block: Block) -> None:
        if block.type == "heading":
            paragraph.text = block.text
            paragraph.font.name = self.theme.heading_font
            paragraph.font.bold = True
            paragraph.font.size = Pt(22 if block.level <= 2 else 18)
            paragraph.font.color.rgb = rgb(self.theme.title)
        elif block.type == "paragraph":
            paragraph.text = block.text
            paragraph.font.name = self.theme.body_font
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = rgb(self.theme.body)
        elif block.type in {"bullets", "numbered"}:
            for index, item in enumerate(block.items):
                target = paragraph if index == 0 else paragraph._parent.add_paragraph()
                target.text = f"- {item}" if block.type == "bullets" else f"{index + 1}. {item}"
                target.font.name = self.theme.body_font
                target.font.size = Pt(18)
                target.font.color.rgb = rgb(self.theme.body)
                target.space_after = Pt(7)
        elif block.type == "code":
            paragraph.text = block.text
            paragraph.font.name = "Menlo"
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = rgb(self.theme.body)

    def add_infographic(self, slide, block: Block) -> None:
        self.add_panel(slide, 5.12, 1.45, 4.18, 5.3)
        self.add_text(slide, str(block.data.get("title", "Highlights")), 5.35, 1.58, 3.9, 0.35, 13, self.theme.title, bold=True, align=PP_ALIGN.CENTER, font=self.theme.heading_font)
        variant = block.data.get("variant")
        if variant == "metrics":
            self.add_metric_cards(slide, block)
        elif variant == "process":
            self.add_process_flow(slide, block)
        elif variant == "highlights":
            self.add_highlight_tiles(slide, block)

    def add_metric_cards(self, slide, block: Block) -> None:
        colors = [self.theme.accent, self.theme.accent_2, self.theme.accent_3, self.theme.accent_4]
        for index, metric in enumerate(block.data.get("metrics", [])[:4]):
            if not isinstance(metric, dict):
                continue
            col, row = index % 2, index // 2
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.35 + col * 2.05), Inches(2.05 + row * 1.75), Inches(1.85), Inches(1.35))
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(colors[index % len(colors)])
            shape.line.fill.background()
            frame = shape.text_frame
            frame.clear()
            value = frame.paragraphs[0]
            value.text = str(metric.get("value", ""))
            value.font.name = self.theme.heading_font
            value.font.size = Pt(24)
            value.font.bold = True
            value.font.color.rgb = rgb(self.theme.on_accent)
            value.alignment = PP_ALIGN.CENTER
            label = frame.add_paragraph()
            label.text = str(metric.get("label", "Metric"))
            label.font.name = self.theme.body_font
            label.font.size = Pt(10)
            label.font.color.rgb = rgb(self.theme.on_accent)
            label.alignment = PP_ALIGN.CENTER

    def add_process_flow(self, slide, block: Block) -> None:
        for index, step in enumerate(block.data.get("steps", [])[:5]):
            shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(5.35), Inches(2.0 + index * 0.8), Inches(3.9), Inches(0.58))
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(self.theme.accent)
            shape.line.fill.background()
            frame = shape.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.text = f"{index + 1}. {step}"
            paragraph.font.name = self.theme.body_font
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = rgb(self.theme.on_accent)
            paragraph.alignment = PP_ALIGN.CENTER

    def add_highlight_tiles(self, slide, block: Block) -> None:
        colors = [tint(self.theme.accent, self.theme.background, 0.82), tint(self.theme.accent_2, self.theme.background, 0.82), tint(self.theme.accent_3, self.theme.background, 0.82), tint(self.theme.accent_4, self.theme.background, 0.82)]
        for index, item in enumerate(block.data.get("items", [])[:4]):
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(2.0 + index * 0.82), Inches(3.9), Inches(0.62))
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(colors[index % len(colors)])
            shape.line.color.rgb = rgb(tint(self.theme.muted, self.theme.background, 0.62))
            frame = shape.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.text = str(item)
            paragraph.font.name = self.theme.body_font
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = rgb(self.theme.title)
            paragraph.alignment = PP_ALIGN.CENTER

    def add_header(self, slide, title: str) -> None:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.12))
        accent.fill.solid()
        accent.fill.fore_color.rgb = rgb(self.theme.accent)
        accent.line.fill.background()
        self.add_text(slide, title, 0.75, 0.42, 8.8, 0.65, 28, self.theme.title, bold=True, font=self.theme.heading_font)

    def add_panel(self, slide, left: float, top: float, width: float, height: float) -> None:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = rgb(self.theme.surface)
        panel.line.color.rgb = rgb(tint(self.theme.background, self.theme.body, 0.1))

    def add_text(self, slide, text: str, left: float, top: float, width: float, height: float, size: int, color: tuple[int, int, int], bold: bool = False, align=None, font: str | None = None) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = font or self.theme.body_font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = rgb(color)
        if align:
            paragraph.alignment = align

    def add_image(self, slide, block: Block, image_base_path: Path) -> None:
        if not block.src:
            return
        image_path = Path(block.src).expanduser()
        if not image_path.is_absolute():
            image_path = image_base_path / image_path
        if image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(2.0), width=Inches(8.0))

    def add_notes(self, slide, notes: str) -> None:
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def set_background(self, slide, color: tuple[int, int, int]) -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(color)


def rgb(color: tuple[int, int, int]):
    return RGBColor(color[0], color[1], color[2])


def tint(foreground: tuple[int, int, int], background: tuple[int, int, int], amount: float) -> tuple[int, int, int]:
    return tuple(int(round(foreground[index] * (1 - amount) + background[index] * amount)) for index in range(3))


def deck_to_dict(deck: Deck) -> dict[str, object]:
    return {
        "title": deck.title,
        "source_path": str(deck.source_path) if deck.source_path else None,
        "metadata": deck.metadata,
        "slides": [
            {
                "title": slide.title,
                "subtitle": slide.subtitle,
                "notes": slide.notes,
                "layout": slide.layout,
                "blocks": [block.__dict__ for block in slide.blocks],
            }
            for slide in deck.slides
        ],
    }


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:
        print(f"error: {exc}", file=sys.stderr)
        raise SystemExit(1)
